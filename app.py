import streamlit as st
import google.generativeai as genai
import pdfplumber
from PIL import Image
import io
import requests
from bs4 import BeautifulSoup
import docx
from streamlit_mic_recorder import mic_recorder
from gtts import gTTS
import base64
import os
import sqlite3
import json

# --- RAG LIBRARIES ---
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain.chains import create_retrieval_chain
from langchain_google_genai import ChatGoogleGenerativeAI

# --- EXPORT LIBRARIES ---
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch

# --- Configuration ---
try:
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
    MODEL_NAME = 'gemini-2.5-flash'
except Exception as e:
    st.error(f"Error configuring Gemini API: {e}")
    st.error("Please make sure you have a valid GEMINI_API_KEY in your .streamlit/secrets.toml file.")
    st.stop()

PROJECTS_DIR = "projects"
os.makedirs(PROJECTS_DIR, exist_ok=True)

# --- DATABASE (Features 2 & 3) ---
def init_db():
    conn = sqlite3.connect('projects.db')
    c = conn.cursor()
    c.execute('''
    CREATE TABLE IF NOT EXISTS projects (
        id INTEGER PRIMARY KEY,
        name TEXT UNIQUE NOT NULL,
        index_path TEXT NOT NULL,
        image_paths TEXT,
        persona TEXT NOT NULL
    )
    ''')
    conn.commit()
    conn.close()

def list_projects():
    conn = sqlite3.connect('projects.db')
    c = conn.cursor()
    c.execute("SELECT name FROM projects")
    projects = [row[0] for row in c.fetchall()]
    conn.close()
    return projects

def create_project(name, index_path, image_paths, persona):
    conn = sqlite3.connect('projects.db')
    c = conn.cursor()
    try:
        c.execute("INSERT INTO projects (name, index_path, image_paths, persona) VALUES (?, ?, ?, ?)",
                  (name, index_path, json.dumps(image_paths), persona))
        conn.commit()
        st.sidebar.success(f"Project '{name}' created!")
    except sqlite3.IntegrityError:
        st.sidebar.error(f"A project with name '{name}' already exists.")
    finally:
        conn.close()

def get_project(name):
    conn = sqlite3.connect('projects.db')
    c = conn.cursor()
    c.execute("SELECT index_path, image_paths, persona FROM projects WHERE name=?", (name,))
    project = c.fetchone()
    conn.close()
    if project:
        return {"index_path": project[0], "image_paths": json.loads(project[1]), "persona": project[2]}
    return None

# --- Custom CSS ---
def load_css():
    st.markdown("""
    <style>
    body { color: #e6e6e6; }
    [data-testid="stAppViewContainer"] { background-color: #0b141a; }
    [data-testid="stSidebar"] { background-color: #1a1f2b; }
    [data-testid="stSidebar"] * { color: #fafafa; }
    [data-testid="stChatMessage"] { border-radius: 12px; padding: 0.9em 1.1em; box-shadow: 0 2px 4px rgba(0,0,0,0.2); border: none; }
    [data-testid="stChatMessage"][data-testid-user="true"] { background-color: #056162; color: #ffffff; border-bottom-right-radius: 4px; }
    [data-testid="stChatMessage"][data-testid-user="false"] { background-color: #262730; color: #fafafa; border-bottom-left-radius: 4px; }
    [data-testid="stButton"] button { border-radius: 8px; }
    h1, h2, h3, h4, h5, h6 { color: #fafafa; }
    .stTabs [data-baseweb="tab"] { color: #a0a0a0; }
    .stTabs [data-baseweb="tab"][aria-selected="true"] { color: #fafafa; }
    [data-testid="stTextInput"] input, [data-testid="stTextArea"] textarea { color: #fafafa; background-color: #262730; border-radius: 20px; }
    [data-testid="stTextInput"] div[data-baseweb="input"] { background-color: #262730; border-radius: 20px; }
    </style>
    """, unsafe_allow_html=True)

# --- Text-to-Speech ---
@st.cache_data(show_spinner="Generating audio...")
def text_to_speech(text):
    audio_fp = io.BytesIO()
    tts = gTTS(text=text, lang='en')
    tts.write_to_fp(audio_fp)
    audio_fp.seek(0)
    return audio_fp.read()

# --- RAG Core Functions ---
def save_uploaded_images(uploaded_images, project_folder):
    image_paths = []
    for img_file in uploaded_images:
        img_path = os.path.join(project_folder, img_file.name)
        with open(img_path, "wb") as f:
            f.write(img_file.getvalue())
        image_paths.append(img_path)
    return image_paths

@st.cache_data(show_spinner="Processing and indexing documents...")
def build_vector_index(uploaded_files_tuple, url_input, project_folder):
    all_texts = []
    preview_parts = []
    
    # ... (File/URL processing logic - no change) ...
    for uploaded_file in uploaded_files_tuple:
        file_bytes = io.BytesIO(uploaded_file.getvalue())
        file_name = uploaded_file.name
        text = ""
        try:
            if file_name.endswith(".pdf"):
                with pdfplumber.open(file_bytes) as pdf:
                    text = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
            elif file_name.endswith(".docx"):
                doc = docx.Document(file_bytes)
                text = "\n".join([p.text for p in doc.paragraphs if p.text])
            elif file_name.endswith(".txt"):
                text = file_bytes.getvalue().decode('utf-8')
            if text:
                all_texts.append(text)
                preview_parts.append({"name": file_name, "type": "text", "content": text})
        except Exception as e:
            st.error(f"Error processing file {file_name}: {e}")
    if url_input:
        try:
            response = requests.get(url_input)
            response.raise_for_status() 
            soup = BeautifulSoup(response.content, 'html.parser')
            text = " ".join([elem.get_text() for elem in soup.find_all(['p', 'h1', 'h2', 'h3', 'li', 'article'])])
            if text:
                all_texts.append(text)
                preview_parts.append({"name": f"Website: {url_input}", "type": "text", "content": text})
        except requests.RequestException as e:
            st.error(f"Error fetching URL: {e}")
    
    if not all_texts:
        return None, preview_parts

    # 3. Split Texts
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    all_chunks = text_splitter.create_documents(all_texts)
    
    # 4. Create Embeddings & Index
    try:
        model_name = "all-MiniLM-L6-v2"
        embeddings = HuggingFaceEmbeddings(model_name=model_name)
        vector_index = FAISS.from_documents(all_chunks, embeddings)
        
        index_path = os.path.join(project_folder, "faiss_index")
        vector_index.save_local(index_path)
        
        return index_path, preview_parts
    except Exception as e:
        st.error(f"Error creating vector index: {e}")
        return None, preview_parts

# --- EXPORT FUNCTIONS (Feature 5) ---

def parse_simple_markdown(text):
    lines = text.split('\n')
    parsed_content = []
    for line in lines:
        if line.startswith('## '):
            parsed_content.append(('h2', line[3:].strip()))
        elif line.startswith('# '):
            parsed_content.append(('h1', line[2:].strip()))
        elif line.startswith('* '):
            parsed_content.append(('bullet', line[2:].strip()))
        elif line.strip(): # Non-empty line
            parsed_content.append(('p', line))
    return parsed_content

def export_to_pptx(text):
    """Exports the last response to a PowerPoint file."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    
    parsed_content = parse_simple_markdown(text)
    
    # Get the content placeholder
    content_placeholder = slide.placeholders[1]
    
    title_set = False
    for type, content in parsed_content:
        if (type == 'h1' or type == 'h2') and not title_set:
            slide.shapes.title.text = content
            title_set = True
        else:
            # Add paragraph to the content placeholder's text frame
            p = content_placeholder.text_frame.add_paragraph()
            p.text = content
            if type == 'bullet':
                p.level = 1
            elif type == 'h1' or type == 'h2':
                p.font.bold = True
                p.font.size = Pt(20)
                
    if not title_set:
        slide.shapes.title.text = "Chat Summary"

    f = io.BytesIO()
    prs.save(f)
    f.seek(0)
    return f.getvalue()

def export_to_docx(text):
    doc = docx.Document()
    parsed_content = parse_simple_markdown(text)
    
    for type, content in parsed_content:
        if type == 'h1':
            doc.add_heading(content, level=1)
        elif type == 'h2':
            doc.add_heading(content, level=2)
        elif type == 'bullet':
            doc.add_paragraph(content, style='List Bullet')
        else:
            doc.add_paragraph(content)

    f = io.BytesIO()
    doc.save(f)
    f.seek(0)
    return f.getvalue()

def export_to_pdf(text):
    """Exports the last response to a PDF file."""
    f = io.BytesIO()
    doc = SimpleDocTemplate(f, pagesize=letter, leftMargin=1*inch, rightMargin=1*inch)
    styles = getSampleStyleSheet()
    story = []
    
    styles.add(ParagraphStyle(
        name='MyH1', 
        parent=styles['Heading1'], 
        fontSize=16, 
        spaceAfter=6
    ))
    styles.add(ParagraphStyle(
        name='MyH2', 
        parent=styles['Heading2'], 
        fontSize=14, 
        spaceAfter=4
    ))
    styles.add(ParagraphStyle(
        name='MyBullet', 
        parent=styles['Bullet'], 
        leftIndent=20
    ))

    parsed_content = parse_simple_markdown(text)
    
    for type, content in parsed_content:
        if type == 'h1':
            story.append(Paragraph(content, styles['MyH1'])) # Use new name
        elif type == 'h2':
            story.append(Paragraph(content, styles['MyH2'])) # Use new name
        elif type == 'bullet':
            story.append(Paragraph(f"• {content}", styles['MyBullet'])) # Use new name
        else:
            story.append(Paragraph(content, styles['BodyText']))
        story.append(Spacer(1, 0.1*inch))
        
    doc.build(story)
    f.seek(0)
    return f.getvalue()
def format_chat_for_export(messages):
    export_string = "Your Chat History\n" + "="*20 + "\n\n"
    for msg in messages:
        role = "You" if msg["role"] == "user" else "Assistant"
        export_string += f"{role}:\n{msg['content']}\n\n---\n\n"
    return export_string

# --- Streamlit App ---

st.set_page_config(page_title="RAG Chat App", layout="wide")
load_css()
init_db()

# --- Session State Initialization ---
if "messages" not in st.session_state:
    st.session_state.messages = []
if "system_instruction" not in st.session_state:
    st.session_state.system_instruction = "" # This will be set by the persona selector
if "preview_content" not in st.session_state:
    st.session_state.preview_content = []
if "prompt_from_button" not in st.session_state:
    st.session_state.prompt_from_button = None
if "retriever" not in st.session_state:
    st.session_state.retriever = None
if "llm" not in st.session_state:
    st.session_state.llm = None
if "image_content" not in st.session_state:
    st.session_state.image_content = []
if "current_project" not in st.session_state:
    st.session_state.current_project = None

# --- Sidebar ---
with st.sidebar:
    st.header("Project Management")

    st.subheader("1. Load Project")
    projects = list_projects()
    if not projects:
        st.caption("No projects found. Create one below.")
    
    selected_project = st.selectbox("Select a project:", projects, index=None, placeholder="Select a project...")
    
    if st.button("Load Project", use_container_width=True):
        if not selected_project:
            st.error("Please select a project to load.")
        else:
            with st.spinner(f"Loading project '{selected_project}'..."):
                project_data = get_project(selected_project)
                if project_data:
                    try:
                        model_name = "all-MiniLM-L6-v2"
                        embeddings = HuggingFaceEmbeddings(model_name=model_name)
                        vector_index = FAISS.load_local(
                            project_data["index_path"], 
                            embeddings,
                            allow_dangerous_deserialization=True 
                        )
                        
                        st.session_state.image_content = []
                        st.session_state.preview_content = []
                        
                        # --- ### FIX 2: Load IMAGE preview content ### ---
                        for img_path in project_data["image_paths"]:
                            try:
                                img = Image.open(img_path)
                                st.session_state.image_content.append(img)
                                st.session_state.preview_content.append({"name": os.path.basename(img_path), "type": "image", "content": img})
                            except FileNotFoundError:
                                st.warning(f"Image file not found: {img_path}")
                        
                        # --- ### FIX 2: Load TEXT preview content ### ---
                        try:
                            # Get all documents from the vector store
                            all_docs = vector_index.docstore._dict
                            # Process them into the preview format
                            for doc_id, doc in all_docs.items():
                                st.session_state.preview_content.append({
                                    "name": doc.metadata.get('source', 'source_doc') + f" (chunk {doc_id})",
                                    "type": "text",
                                    "content": doc.page_content
                                })
                        except Exception as e:
                            st.warning(f"Could not load text previews: {e}")
                        
                        # Load Persona
                        st.session_state.system_instruction = project_data["persona"]
                        
                        st.session_state.llm = ChatGoogleGenerativeAI(model=MODEL_NAME, google_api_key=st.secrets["GEMINI_API_KEY"])
                        st.session_state.retriever = vector_index.as_retriever(search_kwargs={"k": 8})
                        
                        st.session_state.current_project = selected_project
                        st.session_state.messages = []
                        st.success(f"Project '{selected_project}' loaded!")
                        st.rerun()
                        
                    except Exception as e:
                        st.error(f"Error loading project: {e}")
                else:
                    st.error("Project data not found.")

    st.divider()

    st.subheader("2. Create New Project")
    new_project_name = st.text_input("New Project Name:")
    
    personae = {
        "Helpful Assistant": "You are a helpful assistant. Answer the user's questions clearly, using the provided context.",
        "Concise Expert (with Citations)": "You are an expert. Answer with key bullet points. *You MUST cite your sources from the context.* Be brief.",
        "Detailed Explainer (with Citations)": (
            "You are a teacher. Your goal is to explain the answer in great detail. "
            "*First, use the provided context* to form the core of your answer, and you *MUST cite sources* from it. "
            "*Then, expand on the answer* with your own general knowledge to provide a full, comprehensive, multi-paragraph explanation. "
            "Provide step-by-step reasoning."
        ),
        "Tutor": "You are a tutor. Explain the topic step-by-step, as if to a beginner. Use simple language and analogies. *Use the context as your primary source,* but feel free to add simple clarifying details.",
        "Exam Ready": "You are an exam prepper. Be short, crisp, and factual. Use bullet points and bold keywords. *Focus *only on the provided context.** Do not add extra information.",
        "Presentation": "You are a business analyst. Format the answer as 3-5 clear, impactful bullet points for a PowerPoint slide. *Use the context* to find the key data and insights."
    }
    # This selector is now "live"
    st.session_state.system_instruction = st.selectbox(
        "Choose AI Persona:", 
        personae.keys(), 
        key="persona_select",
        # Set index based on loaded persona if it exists
        index=list(personae.keys()).index(st.session_state.system_instruction) if st.session_state.system_instruction in personae else 0
    )
    
    uploaded_files = st.file_uploader("Upload files (PDF, DOCX, TXT)", type=["pdf", "docx", "txt"], accept_multiple_files=True)
    uploaded_images = st.file_uploader("Upload images (PNG, JPG)", type=["png", "jpg", "jpeg"], accept_multiple_files=True)
    url_input = st.text_input("Or enter a website URL", placeholder="https.://")
    
    if st.button("Create & Index Project", use_container_width=True):
        if not new_project_name:
            st.error("Please enter a project name.")
        elif not uploaded_files and not url_input and not uploaded_images:
            st.error("Please add files, images, or a URL to create a project.")
        else:
            with st.spinner(f"Creating project '{new_project_name}'..."):
                project_folder = os.path.join(PROJECTS_DIR, new_project_name.lower().replace(" ", "_"))
                os.makedirs(project_folder, exist_ok=True)
                
                image_paths = save_uploaded_images(uploaded_images, project_folder)
                index_path, preview_parts = build_vector_index(tuple(uploaded_files), url_input, project_folder)
                
                create_project(
                    new_project_name, 
                    index_path, 
                    image_paths, 
                    st.session_state.system_instruction # Save the key name
                )
                
                st.info("Project created. You can now load it from the dropdown above.")

    st.divider()
    
    st.subheader("3. Actions")
    if st.session_state.current_project:
        if st.button("Clear Chat History", use_container_width=True):
            st.session_state.messages = []
            st.info("Chat history cleared.")
            st.rerun()

        if st.session_state.messages:
            st.download_button("Export Full Chat (.txt)", format_chat_for_export(st.session_state.messages), "chat_history.txt", "text/plain", use_container_width=True)
            
            st.markdown("*Export Last Response:*")
            last_response = st.session_state.messages[-1]['content']
            c1, c2, c3 = st.columns(3)
            with c1:
                st.download_button("to DOCX", data=export_to_docx(last_response), file_name="export.docx", use_container_width=True)
            with c2:
                st.download_button("to PPTX", data=export_to_pptx(last_response), file_name="export.pptx", use_container_width=True)
            with c3:
                st.download_button("to PDF", data=export_to_pdf(last_response), file_name="export.pdf", use_container_width=True)

# --- Main UI: Tabs ---
if not st.session_state.current_project:
    st.info("Please create or load a project from the sidebar to begin chatting.")
else:
    st.title(f"Project: {st.session_state.current_project}")
    tab1, tab2 = st.tabs(["💬 Chat", "View Documents"])

    with tab1:
        if not st.session_state.messages:
            st.markdown("#### Quick Actions")
            cols = st.columns(3)
            if cols[0].button("Summarize all content", use_container_width=True):
                st.session_state.prompt_from_button = "Summarize all the provided content."
            if cols[1].button("List 3 key takeaways", use_container_width=True):
                st.session_state.prompt_from_button = "What are the 3 key takeaways from all the documents?"
            if cols[2].button("Generate quiz questions", use_container_width=True):
                st.session_state.prompt_from_button = "Generate 5 quiz questions based on the content."

        for i, message in enumerate(st.session_state.messages):
            with st.chat_message(message["role"]):
                if message["role"] == "user" and message["content"] == "[User sent audio]":
                    st.markdown(message["content"])
                else:
                    st.markdown(message["content"])
                if message["role"] == "assistant":
                    if st.button("Read Aloud", key=f"tts_{i}"):
                        audio_bytes = text_to_speech(message["content"])
                        if audio_bytes:
                            st.audio(audio_bytes, format="audio/mp3")

        prompt_parts = []
        text_prompt = st.chat_input("Ask a question about your project...")
        
        input_cols = st.columns([0.85, 0.15])
        with input_cols[1]:
            audio_data = mic_recorder(start_prompt="🎤", stop_prompt="⏹", key='recorder', use_container_width=True)

        if audio_data:
            prompt_parts.append({"mime_type": "audio/wav", "data": audio_data['bytes']})
            st.session_state.messages.append({"role": "user", "content": "[User sent audio]"})
            
        if text_prompt:
            prompt_parts.append(text_prompt)
            st.session_state.messages.append({"role": "user", "content": text_prompt})

        # --- ### FIX 1: Corrected "Quick Action" Logic ### ---
        if st.session_state.prompt_from_button:
            text_prompt = st.session_state.prompt_from_button
            st.session_state.prompt_from_button = None # Clear it
            
            # Manually add to prompt_parts and history
            if text_prompt not in [msg['content'] for msg in st.session_state.messages]:
                 prompt_parts.append(text_prompt)
                 st.session_state.messages.append({"role": "user", "content": text_prompt})


        if prompt_parts:
            if not st.session_state.retriever and not st.session_state.image_content:
                st.warning("This project seems empty. Please check the project or create a new one.")
            else:
                with st.chat_message("user"):
                    # Display the prompt that was just added to history
                    st.markdown(st.session_state.messages[-1]['content'])
                    if audio_data: st.audio(audio_data['bytes'], format="audio/wav")
                    
                with st.spinner("Thinking..."):
                    try:
                        full_response = ""
                        
                        # Check if this query is multimodal (if project has images OR user just sent audio)
                        is_multimodal_query = bool(st.session_state.image_content or audio_data)
                        
                        # Get the currently selected persona prompt
                        current_persona_prompt = personae[st.session_state.system_instruction]

                        if is_multimodal_query:
                            st.info("Multimodal query detected. Answering based on images/audio, with text context if available.")
                            model_input = st.session_state.image_content + prompt_parts
                            
                            if st.session_state.retriever:
                                # Use text_prompt if available, otherwise summarize
                                text_for_rag = text_prompt if text_prompt else "Describe the audio/images."
                                docs = st.session_state.retriever.invoke(text_for_rag)
                                context_text = "\n\n".join([doc.page_content for doc in docs])
                                model_input = [f"--- RELEVANT TEXT CONTEXT ---\n{context_text}\n--- END CONTEXT ---"] + model_input

                            model = genai.GenerativeModel(MODEL_NAME, system_instruction=current_persona_prompt)
                            chat = model.start_chat(history=[])
                            response = chat.send_message(model_input, stream=True)
                            
                            with st.chat_message("assistant"):
                                placeholder = st.empty()
                                for chunk in response:
                                    if chunk.text:
                                        full_response += chunk.text
                                        placeholder.markdown(full_response + "▌")
                                placeholder.markdown(full_response)
                                
                        elif st.session_state.retriever:
                            # --- Build the RAG chain on-the-fly ---
                            prompt_template = ChatPromptTemplate.from_template(
                                f"{current_persona_prompt}\n\n"
                                "Context: {context}\nQuestion: {input}"
                            )
                            document_chain = create_stuff_documents_chain(st.session_state.llm, prompt_template)
                            retrieval_chain = create_retrieval_chain(st.session_state.retriever, document_chain)
                            
                            # Use text_prompt (which could be from the button or chat_input)
                            text_for_rag = text_prompt if text_prompt else ""
                            response_stream = retrieval_chain.stream({"input": text_for_rag})
                            
                            with st.chat_message("assistant"):
                                placeholder = st.empty()
                                for chunk in response_stream:
                                    if "answer" in chunk:
                                        full_response += chunk["answer"]
                                        placeholder.markdown(full_response + "▌")
                                placeholder.markdown(full_response)
                        
                        else:
                            st.error("No valid chat logic path. Please check your project.")

                        st.session_state.messages.append({"role": "assistant", "content": full_response})

                    except Exception as e:
                        st.error(f"Error generating response: {e}")

    with tab2:
        if not st.session_state.preview_content:
            st.info("This project has no documents or images to preview.")
        else:
            st.header("Project Content Preview")
            for item in st.session_state.preview_content:
                with st.expander(f"View: {item['name']}"):
                    if item['type'] == 'text':
                        st.text_area("Content", item['content'], height=300, disabled=True)
                    elif item['type'] == 'image':
                        st.image(item['content'], use_column_width=True)
